import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def add_title_slide(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle is not None:
        slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[tuple[str, int]]):
    """Add a Title and Content slide with bullets.

    bullets: list of (text, level) where level is 0 for top-level, 1 for nested, etc.
    """
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = text
        p.level = max(0, level)
        p.font.size = Pt(20)


def add_section_header(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    slide_layout = prs.slide_layouts[2]  # Section Header
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle is not None:
        ph = slide.placeholders[1]
        ph.text = subtitle
        for p in ph.text_frame.paragraphs:
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def build_deck(output_path: str) -> None:
    prs = Presentation()

    # Cover
    add_title_slide(
        prs,
        title="Intelligent Banking Operations Agent",
        subtitle="Fraud & Credit Triage – Hackathon Pitch",
    )

    # Agenda (simplified)
    add_bullets_slide(
        prs,
        title="Agenda",
        bullets=[
            ("The problem we solve", 0),
            ("Our solution: Fraud & Credit triage", 0),
            ("Impact: faster decisions, fewer losses", 0),
            ("Results & next steps", 0),
        ],
    )

    # Fraud section
    add_section_header(prs, "Fraud Triage", "Minimize losses with low friction and efficient reviews")

    add_bullets_slide(
        prs,
        "Goal",
        bullets=[
            ("Prioritize suspicious activity to cut fraud losses", 0),
            ("Control review costs and customer friction", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "How it works (simple)",
        bullets=[
            ("We score each event for risk", 0),
            ("High risk → quick hold; medium → review; low → pass", 0),
            ("Analysts see why it was flagged and decide fast", 0),
            ("The system learns from decisions to improve precision", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Why it works",
        bullets=[
            ("Known checks catch obvious fraud; AI finds new patterns", 0),
            ("Clear reasons shown to analysts → trust + speed", 0),
            ("Right alerts at the right time → less noise", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "High-signal features",
        bullets=[
            ("Velocity: txn count/amount over 1h/24h/7d; rate-of-change vs baseline", 0),
            ("Geo: distance from last location, impossible travel, country risk", 0),
            ("Device/IP: new vs seen, fingerprint diversity, proxy/ASN risk", 0),
            ("Behavioral: z-score vs user history, time-of-day, MCC mix", 0),
            ("Merchant: MCC risk, merchant novelty, basket deviation", 0),
            ("Graph: shared devices/IPs across accounts, degree features", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Model choices",
        bullets=[
            ("Isolation Forest (unsupervised): cold start, imbalance, ms latency", 0),
            ("Autoencoder/LOF: alternate anomaly options", 0),
            ("XGBoost/Logistic (supervised) when labels exist; support calibration", 0),
            ("Calibration: map scores to risk buckets/empirical fraud rates", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Business impact",
        bullets=[
            ("Fewer false alarms → happier customers", 0),
            ("Faster decisions → more throughput", 0),
            ("Less fraud loss at the same review effort", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "What we measure",
        bullets=[
            ("Accuracy of alerts", 0),
            ("Time to decision", 0),
            ("Estimated loss avoided", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Analyst experience",
        bullets=[
            ("Clear 'why flagged' explanation", 0),
            ("Similar past cases for quick context", 0),
            ("One-click actions (hold/release/report)", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Handling realities",
        bullets=[
            ("Imbalance & delayed labels: anomaly/PU; evaluate via PR curves", 0),
            ("Drift: weekly PSI; recalibrate thresholds; review spiking rules", 0),
            ("Adversaries: multi-signal, random spot-checks, rotating windows", 0),
            ("Data quality: robust defaults; schema validation; fallback rules", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Operations",
        bullets=[
            ("Real-time friendly", 0),
            ("Scales with demand", 0),
            ("Safe fallback if AI is unavailable", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Governance & pitfalls",
        bullets=[
            ("Version rules/models; immutable logs; approvals; PII minimization", 0),
            ("Avoid: rule explosion, ROC over PR, ignoring capacity, static thresholds", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Fraud triage – 30s story",
        bullets=[
            ("We combine proven checks with AI to spot fraud early", 0),
            ("We route only meaningful alerts to the team", 0),
            ("Analysts act fast with clear reasons", 0),
        ],
    )

    # Credit section
    add_section_header(prs, "Credit Triage", "Approve good customers fast; control portfolio risk")

    add_bullets_slide(
        prs,
        "Goal",
        bullets=[
            ("Maximize portfolio profit under risk, capacity, and compliance", 0),
            ("Approve good, block risky, refer ambiguous", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "How it works (simple)",
        bullets=[
            ("We estimate risk and affordability", 0),
            ("Clear decisions: approve, review, or decline", 0),
            ("Approved cases get limits/pricing that fit ability to pay", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Why it works",
        bullets=[
            ("Rules keep us compliant; AI improves approvals with low risk", 0),
            ("Reason codes explain decisions to customers and regulators", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "High-signal features",
        bullets=[
            ("Affordability: DTI/FOIR, income stability, cash flow, EMIs", 0),
            ("Credit behavior: bureau score, delinquencies, utilization, inquiries", 0),
            ("Stability/identity: employment tenure, address/device consistency", 0),
            ("Fraud overlap: synthetic markers, app velocity, shared device/IP", 0),
            ("Macro/segment: industry, region proxy, secured vs unsecured", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "What we deliver",
        bullets=[
            ("Higher approval rate at the same bad rate", 0),
            ("Fair, explainable decisions with audit trail", 0),
            ("Right limits and pricing from day one", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Thresholds, queues, actions",
        bullets=[
            ("Approve if PD < t1; decline if PD > t2; else Refer", 0),
            ("Maximize Expected Profit ≈ margin*EAD*(1−PD) − LGD*EAD*PD − OpCost", 0),
            ("Capacity-aware Refer volume; segment by product/bureau band", 0),
            ("On Approve: risk-based pricing and affordability caps", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "What we measure",
        bullets=[
            ("Approval rate and default rate", 0),
            ("Time to decision", 0),
            ("Customer satisfaction and fairness checks", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Explainability & adverse action",
        bullets=[
            ("Reason codes mapping to top contributors", 0),
            ("What-if: path from Refer → Approve", 0),
            ("Audit trail: inputs, score, policy hits, versions, final decision", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Handling realities",
        bullets=[
            ("Selection bias: parceling/EM; validate via sensitivity/shadow cohorts", 0),
            ("Macro shocks: adjust cutoffs/limits; faster vintages; stress tests", 0),
            ("Data quality: verify income; robust imputation; conservative fallback", 0),
            ("Fair lending: disparity monitoring; reason-code integrity; consistent cutoffs", 0),
            ("Adversaries: inflated income/docs → cross-check with bank data", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Operations & trust",
        bullets=[
            ("Fast pre-qualification; quick underwriting", 0),
            ("Built-in A/B testing for safe improvements", 0),
            ("Compliant, auditable decisions", 0),
        ],
    )

    add_bullets_slide(
        prs,
        "Credit triage – 30s story",
        bullets=[
            ("We approve good customers quickly and fairly", 0),
            ("Risk and affordability guide limits and pricing", 0),
            ("Clear reasons build trust and compliance", 0),
        ],
    )

    # Results summary (executive)
    add_bullets_slide(
        prs,
        "Results summary",
        bullets=[
            ("Fraud: fewer false alerts and faster reviews", 0),
            ("Credit: higher approvals at stable risk", 0),
            ("Ops: quicker decisions and clear reasons", 0),
            ("Governance: fully auditable and compliant", 0),
        ],
    )

    # Closing
    add_section_header(prs, "Q&A", "Thank you")

    prs.save(output_path)


if __name__ == "__main__":
    out = os.path.abspath("Hackathon_Fraud_Credit_Triage.pptx")
    build_deck(out)
    print(f"Generated deck: {out}")



